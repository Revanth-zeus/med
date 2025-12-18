import io
import re
from typing import List, Dict
import PyPDF2
from docx import Document
from pptx import Presentation

class DocumentProcessor:
    """Process PDF, DOCX, and PPTX files - OPTIMIZED for any size"""
    
    @staticmethod
    def process_pdf(file_bytes: bytes) -> Dict:
        """Extract text from PDF - fast extraction"""
        try:
            pdf_file = io.BytesIO(file_bytes)
            pdf_reader = PyPDF2.PdfReader(pdf_file)
            
            text_chunks = []
            full_text = ""
            total_pages = len(pdf_reader.pages)
            
            print(f"   📖 Extracting {total_pages} pages...")
            
            for page_num, page in enumerate(pdf_reader.pages, 1):
                page_text = page.extract_text()
                full_text += page_text + "\n\n"
                
                # Chunk by page for better context
                if page_text.strip():
                    text_chunks.append({
                        "page": page_num,
                        "text": page_text.strip()
                    })
                
                # Progress every 10 pages
                if page_num % 10 == 0:
                    print(f"   📄 Extracted {page_num}/{total_pages} pages...")
            
            return {
                "full_text": full_text,
                "chunks": text_chunks,
                "total_pages": total_pages
            }
        
        except Exception as e:
            print(f"Error processing PDF: {e}")
            return {"full_text": "", "chunks": [], "total_pages": 0}
    
    @staticmethod
    def process_docx(file_bytes: bytes) -> Dict:
        """Extract text from DOCX - fast extraction"""
        try:
            docx_file = io.BytesIO(file_bytes)
            doc = Document(docx_file)
            
            text_chunks = []
            full_text = ""
            
            # Extract paragraphs
            for para_num, para in enumerate(doc.paragraphs, 1):
                if para.text.strip():
                    full_text += para.text + "\n"
                    text_chunks.append({
                        "paragraph": para_num,
                        "text": para.text.strip()
                    })
            
            # Extract tables
            for table_num, table in enumerate(doc.tables, 1):
                table_text = "\n".join([
                    " | ".join([cell.text for cell in row.cells])
                    for row in table.rows
                ])
                if table_text.strip():
                    full_text += f"\nTable {table_num}:\n{table_text}\n"
                    text_chunks.append({
                        "table": table_num,
                        "text": table_text
                    })
            
            return {
                "full_text": full_text,
                "chunks": text_chunks,
                "total_paragraphs": len(doc.paragraphs)
            }
        
        except Exception as e:
            print(f"Error processing DOCX: {e}")
            return {"full_text": "", "chunks": [], "total_paragraphs": 0}
    
    @staticmethod
    def process_pptx(file_bytes: bytes) -> Dict:
        """Extract text from PPTX"""
        try:
            pptx_file = io.BytesIO(file_bytes)
            prs = Presentation(pptx_file)
            
            text_chunks = []
            full_text = ""
            
            for slide_num, slide in enumerate(prs.slides, 1):
                slide_text = []
                
                # Extract text from shapes
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        slide_text.append(shape.text)
                
                if slide_text:
                    combined_text = "\n".join(slide_text)
                    full_text += f"Slide {slide_num}:\n{combined_text}\n\n"
                    text_chunks.append({
                        "slide": slide_num,
                        "text": combined_text
                    })
            
            return {
                "full_text": full_text,
                "chunks": text_chunks,
                "total_slides": len(prs.slides)
            }
        
        except Exception as e:
            print(f"Error processing PPTX: {e}")
            return {"full_text": "", "chunks": [], "total_slides": 0}
    
    @staticmethod
    def smart_chunk_text(text: str, chunk_size: int = 1500, overlap: int = 150) -> List[str]:
        """
        Smart text chunking - OPTIMIZED for speed, NO LIMITS
        
        Args:
            text: Full text to chunk
            chunk_size: Target size for each chunk (increased to 1500 for fewer chunks)
            overlap: Overlap between chunks
        
        Returns:
            List of text chunks
        """
        # Clean text
        text = re.sub(r'\s+', ' ', text).strip()
        
        if len(text) <= chunk_size:
            return [text] if text else []
        
        chunks = []
        start = 0
        text_len = len(text)
        
        while start < text_len:
            end = start + chunk_size
            
            # Try to break at sentence boundary
            if end < text_len:
                # Look for sentence endings near the chunk boundary
                sentence_end = text.rfind('.', start, end + 100)
                if sentence_end > start + chunk_size // 2:  # Only if we found one in reasonable range
                    end = sentence_end + 1
            
            chunk = text[start:end].strip()
            if chunk:
                chunks.append(chunk)
            
            start = end - overlap
            if start >= text_len:
                break
        
        # NO LIMIT - process all chunks
        print(f"   📦 Created {len(chunks)} chunks from {text_len:,} characters")
        
        return chunks
    
    @classmethod
    def process_document(cls, file_bytes: bytes, mime_type: str, filename: str) -> Dict:
        """
        Process document based on MIME type - FAST, NO LIMITS
        
        Returns:
            Dict with processed content and metadata
        """
        result = {
            "filename": filename,
            "mime_type": mime_type,
            "full_text": "",
            "chunks": [],
            "metadata": {}
        }
        
        if 'pdf' in mime_type:
            processed = cls.process_pdf(file_bytes)
            result["full_text"] = processed["full_text"]
            result["chunks"] = [chunk["text"] for chunk in processed["chunks"]]
            result["metadata"]["total_pages"] = processed["total_pages"]
        
        elif 'wordprocessingml' in mime_type or 'msword' in mime_type:
            processed = cls.process_docx(file_bytes)
            result["full_text"] = processed["full_text"]
            result["chunks"] = [chunk["text"] for chunk in processed["chunks"]]
            result["metadata"]["total_paragraphs"] = processed["total_paragraphs"]
        
        elif 'presentationml' in mime_type or 'ms-powerpoint' in mime_type:
            processed = cls.process_pptx(file_bytes)
            result["full_text"] = processed["full_text"]
            result["chunks"] = [chunk["text"] for chunk in processed["chunks"]]
            result["metadata"]["total_slides"] = processed["total_slides"]
        
        # Apply smart chunking for better RAG - NO LIMITS
        if result["full_text"]:
            result["smart_chunks"] = cls.smart_chunk_text(result["full_text"])
        
        return result